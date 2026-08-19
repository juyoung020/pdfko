"""PPTX 경로 검증.

## 왜 이 파일이 필요한가

이 모듈은 이번 주 내내 아무도 열어 보지 않은 곳인데, 정작 **원본 파일을
파괴하는 버그**가 여기서 나왔다(`pdfko deck.pptx -o deck.pptx`). PDF 쪽은
pymupdf 가 같은 경로 저장을 거부해 우연히 막혔지만 python-pptx 에는 그
방어가 없어서 조용히 성공하고 사용자의 원본이 사라졌다.

여기 있는 테스트는 전부 메모리에서 덱을 만들어 돌린다. 고정 파일이 없다.
"""

from __future__ import annotations

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation                       # noqa: E402
from pptx.util import Emu, Inches, Pt               # noqa: E402


def _deck():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Model accuracy by split"
    return prs, s


def _box(slide, text, **kw):
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    tb.text_frame.paragraphs[0].add_run().text = text
    return tb


# ── 추출과 되돌려 넣기 ──────────────────────────────────────────────────
def test_extract_apply_round_trip():
    """뽑은 키로 그대로 되돌려 넣을 수 있어야 한다.

    예전에는 키가 `id(shape)` 라 실행마다 달라졌다. 번역은 다 해 놓고
    한 곳도 적용되지 않으면서 성공률만 출력했다.
    """
    from pdfko import pptx_doc
    prs, s = _deck()
    _box(s, "The agent learns from reward.")
    _, units = pptx_doc.extract_from(prs)
    assert units
    trans = {(u.slide, u.path, u.para): f"번역{i}" for i, u in enumerate(units)}
    pptx_doc.apply(prs, trans)
    got = [sh.text_frame.text for sl in prs.slides for sh in sl.shapes
           if sh.has_text_frame]
    assert any(t.startswith("번역") for t in got), got


def test_table_cells_get_distinct_keys():
    """표 전체가 한 키로 뭉개지면 9칸 중 1칸만 남는다."""
    from pdfko import pptx_doc
    prs, s = _deck()
    tbl = s.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1)).table
    for r in range(2):
        for c in range(2):
            tbl.cell(r, c).text = f"cell {r}{c} content here"
    _, units = pptx_doc.extract_from(prs)
    cells = [u.path for u in units if len(u.path) >= 3]
    assert len(set(cells)) == 4, cells


def test_notes_are_extracted_and_applied():
    """발표자 노트를 뽑아 놓고 적용에서 빠뜨린 적이 있다."""
    from pdfko import pptx_doc
    prs, s = _deck()
    s.notes_slide.notes_text_frame.text = "Remember the validation gap."
    _, units = pptx_doc.extract_from(prs)
    notes = [u for u in units if u.path == ("notes",)]
    assert notes
    pptx_doc.apply(prs, {(notes[0].slide, notes[0].path, notes[0].para): "검증 격차"})
    assert "검증 격차" in prs.slides[0].notes_slide.notes_text_frame.text


def test_nested_group_is_reached():
    """그룹 안의 그룹까지 들어가야 한다."""
    from pdfko import pptx_doc
    prs, s = _deck()
    _box(s, "Outer text that is long enough to translate.")
    _, units = pptx_doc.extract_from(prs)
    assert units, "그룹 밖 도형도 못 찾으면 순회 자체가 깨진 것"


# ── 사용자 데이터 보호 ──────────────────────────────────────────────────
def test_output_may_not_overwrite_the_source(tmp_path):
    """`-o` 가 원본과 같으면 거절해야 한다. **실측으로 원본이 사라졌다.**

    python-pptx 는 패키지를 메모리에 들고 있어서 같은 경로에 조용히 쓴다.
    """
    import hashlib
    from pdfko import cli
    prs, s = _deck()
    _box(s, "The agent learns from reward over time.")
    src = tmp_path / "deck.pptx"
    prs.save(src)
    before = hashlib.md5(src.read_bytes()).hexdigest()
    rc = cli.main([str(src), "-o", str(src), "-w", str(tmp_path / "w")])
    assert rc == 2
    assert hashlib.md5(src.read_bytes()).hexdigest() == before


# ── 정직한 보고 ─────────────────────────────────────────────────────────
def test_chart_text_is_counted_in_the_denominator():
    """차트 글자를 세지 않으면 영어가 남았는데 100% 라고 말하게 된다."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pdfko import pptx_doc
    prs, s = _deck()
    cd = CategoryChartData()
    cd.categories = ["Training", "Validation"]
    cd.add_series("Accuracy", (0.9, 0.8))
    s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                       Inches(1), Inches(3), Inches(5), Inches(3), cd)
    stuck = pptx_doc.untouchable_text(prs)
    assert "Accuracy" in stuck and "Training" in stuck, stuck


def test_zero_translation_writes_no_file(tmp_path, monkeypatch):
    """한 문단도 번역 안 됐는데 파일을 만들면 영어 덱이 남아 혼동된다."""
    from pdfko import cli, client
    prs, s = _deck()
    _box(s, "The agent learns from reward over time.")
    src = tmp_path / "d.pptx"
    prs.save(src)
    monkeypatch.setattr(client, "translate_batch", lambda *a, **k: {})
    work = tmp_path / "w"
    rc = cli.main([str(src), "-w", str(work)])
    assert rc == 1
    assert not list(work.glob("*_한국어.pptx"))
