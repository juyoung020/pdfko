"""PPTX 번역 — 원본 서식을 그대로 두고 글자만 한국어로 바꾼다.

## PDF 와 무엇이 다른가

PDF 는 조판이 이미 확정된 결과물이라, 글자를 바꾸면 자리가 모자라 겹친다.
PPTX 는 **편집 가능한 문서**다. 도형과 텍스트 상자가 객체로 남아 있어서
글자를 바꿔 넣으면 파워포인트가 알아서 다시 흘려 준다.

그래서 접근이 다르다.

  PDF  : 조판을 건드리지 않으려고 애쓴다 → 자리가 없으면 깨진다
  PPTX : 런(run) 단위로 글자만 교체한다 → 서식·애니메이션·위치가 그대로 남는다

## 넘침을 어떻게 다루나

파워포인트는 상자보다 글이 길면 밖으로 흘리거나 잘라 버린다. 자동 맞춤이
켜져 있으면 글씨를 줄여 주지만 꺼져 있는 상자도 많다.

이 모듈은 번역 후 **예상 길이를 재서** 넘칠 것 같으면 글씨를 한 단계
(15%) 줄이고 자동 줄바꿈을 켠다. 8pt 아래로는 내려가지 않는다.
크기가 레이아웃에서 상속된 런은 실제 크기를 알 수 없어 건드리지 않는다.
원문을 버리는 일은 없다 — PPTX 는 PDF 와 달리 되돌릴 필요가 없다.

## 서식 보존

한 문단이 여러 런으로 쪼개져 있는 일이 흔하다(굵게·색깔 때문에).
런마다 따로 번역하면 문장이 토막 나므로 문단 단위로 번역한 뒤 **가장 긴 런**
에 넣는다. 그 런의 서식이 문단 전체에 적용되므로, **문장 속 부분 강조는
보존되지 않는다.** 몇 개를 잃었는지는 보고서에 남는다.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


@dataclass
class TextUnit:
    """번역 대상 하나. 슬라이드·도형·문단을 특정한다.

    `path` 는 도형의 **위치**다 — (그룹 인덱스…, 도형 인덱스[, 행, 열]).
    한때 `id(shape)` 를 썼는데, python-pptx 는 도형을 훑을 때마다 새 프록시
    객체를 만들기 때문에 추출 때와 적용 때의 id 가 절대 같지 않다. 그래서
    번역이 하나도 적용되지 않으면서 성공률만 출력하는 버그가 있었다.
    """
    slide: int
    path: tuple
    para: int
    text: str
    kind: str = "body"        # body | table | notes


@dataclass
class SlideReport:
    slide: int
    units: int = 0
    translated: int = 0
    shrunk: list[str] = field(default_factory=list)
    skipped: list[str] = field(default_factory=list)


def _iter_frames(slide):
    """텍스트를 가진 모든 것을 훑는다 — 도형, 표 칸, 그룹 안쪽, 발표자 노트.

    (경로, 도형, 텍스트프레임, 종류) 를 낸다. 경로는 추출과 적용에서 동일하게
    재현되는 유일한 식별자다. 표는 칸마다 다른 경로를 줘야 한다 — 예전에는
    표 전체가 한 키로 뭉개져 9칸 중 1칸만 남았다.
    """
    def walk(shapes, prefix):
        for i, sh in enumerate(shapes):
            # shape_type 은 분류 불가 도형에서 NotImplementedError 를 던진다.
            # (Keynote·구글슬라이드 내보내기에 흔하다) 그룹 판정에는 불필요하다.
            if hasattr(sh, "shapes"):    # 그룹
                yield from walk(sh.shapes, prefix + (i,))
                continue
            if sh.has_text_frame:
                yield prefix + (i,), sh, sh.text_frame, "body"
            if getattr(sh, "has_table", False):
                for r, row in enumerate(sh.table.rows):
                    for c, cell in enumerate(row.cells):
                        yield prefix + (i, r, c), sh, cell.text_frame, "table"
    yield from walk(slide.shapes, ())
    # 발표자 노트도 같은 순회에 넣는다. 예전에는 추출만 하고 적용에서 빠져서
    # 번역해 놓고 버렸다.
    if slide.has_notes_slide:
        yield ("notes",), None, slide.notes_slide.notes_text_frame, "notes"


_A_T = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
# 차트의 축 이름·계열 이름·범례는 `<a:t>` 가 아니라 문자열 캐시 `<c:v>` 에 산다.
# `<a:t>` 만 훑으면 차트가 통째로 안 보인다 — 실측으로 0개가 나왔다.
_C_V = "{http://schemas.openxmlformats.org/drawingml/2006/chart}v"


def untouchable_text(prs: Presentation) -> list[str]:
    """번역이 **닿지 못하는** 텍스트를 찾아 돌려준다.

    차트(축 제목·계열 이름·범례)와 SmartArt 는 슬라이드가 아니라 별도의 XML
    파트에 글자를 담는다. `slide.shapes` 순회로는 절대 보이지 않는다.

    그런데 그냥 못 본 척하면 **"번역 20/20 (100%)"라고 보고하면서 화면에는
    영어 축 제목이 그대로 남는다.** 도구가 거짓말을 하는 셈이다. 고칠 수
    없다면 최소한 세어서 알려야 한다.
    """
    from xml.etree import ElementTree as ET

    found: list[str] = []
    try:
        parts = prs.part.package.iter_parts()
    except Exception:
        return found
    for part in parts:
        ct = getattr(part, "content_type", "") or ""
        if "chart" not in ct and "diagramData" not in ct:
            continue
        try:
            root = ET.fromstring(part.blob)
        except Exception:
            continue
        for tag in (_A_T, _C_V):
            for el in root.iter(tag):
                s = (el.text or "").strip()
                # 숫자·기호만 있는 칸(데이터 값)은 번역 대상이 아니다
                if s and any(c.isalpha() for c in s) and s not in found:
                    found.append(s)
    return found


def extract(path: str | Path) -> tuple[Presentation, list[TextUnit]]:
    """파일에서 번역할 문단을 모두 뽑아낸다."""
    return extract_from(Presentation(str(path)))


def extract_from(prs: Presentation) -> tuple[Presentation, list[TextUnit]]:
    """이미 열어 둔 덱에서 뽑는다. 파일이 없어도 되므로 시험하기 쉽다."""
    units: list[TextUnit] = []
    for si, slide in enumerate(prs.slides, 1):
        for path, sh, tf, kind in _iter_frames(slide):
            for pi, para in enumerate(tf.paragraphs):
                txt = "".join(r.text for r in para.runs)
                if txt.strip() and any(c.isalpha() for c in txt):
                    units.append(TextUnit(si, path, pi, txt, kind))
    return prs, units


def _fits(text: str, tf, shape, kind: str = "body") -> bool:
    """대략적인 넘침 판정.

    글자 폭(한글 1em, 라틴 0.5em)으로 필요한 줄 수를 어림하고 상자 높이와 견준다.
    파워포인트의 실제 조판을 흉내 낼 수는 없으니 보수적으로만 본다.
    """
    # 표 칸은 도형(표 전체) 크기를 쓰면 뜻이 없다. 칸 크기를 알 수 없으므로
    # 판정을 보류한다 — 파워포인트가 행 높이를 알아서 늘린다.
    if kind == "table" or shape is None:
        return True
    try:
        w = shape.width / 914400 * 72        # EMU → pt
        h = shape.height / 914400 * 72
    except Exception:
        return True
    size = 18.0
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size:
                size = r.font.size.pt
                break
        break
    if w <= 0 or h <= 0:
        return True
    em = sum(1.0 if ("가" <= c <= "힣") else 0.5 for c in text)
    per_line = max(1.0, w / size)    # em 폭 기준 — 보정 계수를 곱하면 안 된다
    lines = em / per_line
    return lines * size * 1.25 <= h * 1.05


def apply(prs: Presentation, translations: dict[tuple, str],
          autofit: bool = True) -> list[SlideReport]:
    """번역문을 되돌려 넣는다.

    서식 보존: 한 문단이 여러 런으로 쪼개져 있으면(굵게·색깔 때문에) 번역문을
    **가장 긴 런**에 넣는다. 예전에는 무조건 첫 런에 넣었는데, `**주의:** 본문`
    같은 문단에서는 첫 런이 굵은 라벨이라 문장 전체가 굵어졌고, 반대로 첫 런이
    장식 없는 조각이면 강조가 통째로 사라졌다. 가장 긴 런이 본문일 확률이 높다.

    여러 런을 하나로 합치는 이상 **문장 속 부분 강조는 보존되지 않는다.**
    보고서에 몇 개를 잃었는지 남긴다.
    """
    reports: list[SlideReport] = []
    for si, slide in enumerate(prs.slides, 1):
        rep = SlideReport(slide=si)
        for path, sh, tf, kind in _iter_frames(slide):
            changed = False
            for pi, para in enumerate(tf.paragraphs):
                ko = translations.get((si, path, pi))
                rep.units += 1
                if not ko or not para.runs:
                    continue
                runs = para.runs
                # 본문일 가능성이 가장 높은 런을 고른다
                main = max(range(len(runs)), key=lambda i: len(runs[i].text))
                if len(runs) > 1:
                    rep.skipped.append(f"s{si}{path}p{pi}: 런 {len(runs)}→1")
                runs[main].text = ko
                for i, r in enumerate(runs):
                    if i != main:
                        r.text = ""
                rep.translated += 1
                changed = True

            # 손대지 않은 상자는 건드리지 않는다. 예전에는 번역이 0건인
            # 상자의 글씨까지 줄여서 멀쩡한 서식을 망가뜨렸다.
            if autofit and changed and sh is not None:
                txt = "".join(p.text for p in tf.paragraphs)
                if not _fits(txt, tf, sh, kind):
                    if _shrink(tf):
                        rep.shrunk.append(str(getattr(sh, "name", "shape")))
        reports.append(rep)
    return reports


def _shrink(tf) -> bool:
    """글씨를 한 단계 줄인다. 크기가 상속된 런은 건드리지 않는다.

    `r.font.size` 가 None 이면 레이아웃에서 물려받는다는 뜻이고, 실제 크기를
    모르는 채 18pt 로 가정해 줄이면 44pt 제목을 15pt 로 만들어 버린다.
    """
    done = False
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size is None:
                continue
            r.font.size = Pt(max(8.0, r.font.size.pt * 0.85))
            done = True
    if done and tf.word_wrap is None:
        # 저자가 명시적으로 끈 상자는 건드리지 않는다. 켜 버리면 가로로
        # 흘러 나가던 한 줄이 세로로 쌓여 슬라이드 밖으로 나간다.
        tf.word_wrap = True
    return done


def save(prs: Presentation, out: str | Path) -> None:
    prs.save(str(out))
